import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime

def create_presentation(json_file='slides.json', output_file_base='nano_banana_presentation', skip_images=True):
    try:
        with open(json_file, 'r', encoding='utf-8') as f:
            data = json.load(f)
            # Handle both formats: direct array or object with 'slides' key
            if isinstance(data, dict) and 'slides' in data:
                slides_data = data['slides']
            elif isinstance(data, list):
                slides_data = data
            else:
                print(f"Error: Unexpected JSON format in {json_file}")
                return
    except FileNotFoundError:
        print(f"Error: {json_file} not found.")
        return

    # Generate a timestamp for the output file
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_file = f"{output_file_base}_{timestamp}.pptx"

    prs = Presentation()
    # Set slide dimensions to 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, slide_data in enumerate(slides_data):
        # Use a blank layout for custom positioning
        slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_left = Inches(0.5)
        title_top = Inches(0.3)
        title_width = Inches(12)
        title_height = Inches(1.0)
        
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_tf = title_box.text_frame
        title_tf.text = slide_data.get('title', 'No Title')
        title_tf.paragraphs[0].font.size = Pt(40)
        title_tf.paragraphs[0].font.bold = True

        # Content (Text) - Left side
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(6.0)
        height = Inches(5.0)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        # Handle content as either string or list
        content = slide_data.get('content', '')
        if isinstance(content, list):
            # Convert list to bullet points
            content_text = '\n'.join([f"• {item}" for item in content])
        else:
            content_text = content
            
        tf.text = content_text
        tf.word_wrap = True
        
        # Adjust font size for content
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(20)

        # Image placeholder - Right side
        img_prompt = slide_data.get('image_prompt', '')
        
        # Always create placeholder for now
        placeholder_left = Inches(7.0)
        placeholder_top = Inches(1.5)
        placeholder_width = Inches(5.8)
        placeholder_height = Inches(4.0)
        
        shape = slide.shapes.add_shape(
            1, # msoShapeRectangle
            placeholder_left, placeholder_top, placeholder_width, placeholder_height
        )
        # Make it look like a placeholder (light gray fill, border)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(238, 238, 238) # Light gray (0xEEEEEE)
        shape.line.color.rgb = RGBColor(136, 136, 136) # Darker gray border (0x888888)
        
        # Add text to the placeholder
        p_tf = shape.text_frame
        p_tf.text = f"Image Placeholder\n\nPrompt:\n{img_prompt}" if img_prompt else "No Image Prompt"
        p_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p_tf.paragraphs[0].font.bold = True
        p_tf.paragraphs[0].font.size = Pt(14)
        
        # Add notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = f"Image Prompt: {img_prompt}"

    prs.save(output_file)
    print(f"✅ Presentation saved to {output_file}")
    print(f"📊 Created {len(slides_data)} slides")

if __name__ == "__main__":
    create_presentation()
