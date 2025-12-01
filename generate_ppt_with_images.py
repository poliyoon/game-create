import json
import os
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime
import base64

def get_api_key():
    try:
        with open('.env', 'r', encoding='utf-8') as f:
            content = f.read().strip()
            # Get the last non-empty line
            lines = [line.strip() for line in content.split('\n') if line.strip() and not line.strip().startswith('#')]
            return lines[-1] if lines else None
    except FileNotFoundError:
        print("Error: .env file not found.")
        return None

def generate_image_with_imagen(prompt, output_path, api_key):
    """Generate image using Imagen 4.0 API"""
    try:
        genai.configure(api_key=api_key)
        
        # Use Imagen 4.0 Fast model for faster generation
        imagen = genai.ImageGenerationModel("imagen-4.0-fast-generate-001")
        
        print(f"  Generating image: {prompt[:50]}...")
        
        # Generate image
        result = imagen.generate_images(
            prompt=prompt,
            number_of_images=1,
            aspect_ratio="16:9",  # Good for presentations
            safety_filter_level="block_some",
            person_generation="allow_adult"
        )
        
        if result.images:
            # Save the first image
            image_data = result.images[0]._image_bytes
            with open(output_path, 'wb') as f:
                f.write(image_data)
            print(f"  ✅ Image saved to {output_path}")
            return True
        else:
            print(f"  ❌ No image generated")
            return False
            
    except Exception as e:
        print(f"  ❌ Failed to generate image: {e}")
        return False

def create_presentation(json_file='slides.json', output_file_base='nano_banana_presentation', generate_images=True):
    api_key = None
    if generate_images:
        api_key = get_api_key()
        if not api_key:
            print("⚠️  No API key found. Will create placeholders instead.")
            generate_images = False

    try:
        with open(json_file, 'r', encoding='utf-8') as f:
            data = json.load(f)
            if isinstance(data, dict) and 'slides' in data:
                slides_data = data['slides']
            elif isinstance(data, list):
                slides_data = data
            else:
                print(f"Error: Unexpected JSON format in {json_file}")
                return
    except FileNotFoundError:
        print(f"Error: {json_file} not found.")
        return

    # Generate a timestamp for the output file
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_file = f"{output_file_base}_{timestamp}.pptx"

    # Create images directory if it doesn't exist
    if generate_images and not os.path.exists('generated_images'):
        os.makedirs('generated_images')

    prs = Presentation()
    # Set slide dimensions to 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    print(f"\n🎨 Creating presentation with {len(slides_data)} slides...\n")

    for i, slide_data in enumerate(slides_data):
        print(f"📄 Slide {i+1}/{len(slides_data)}: {slide_data.get('title', 'No Title')[:50]}...")
        
        # Use a blank layout for custom positioning
        slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_left = Inches(0.5)
        title_top = Inches(0.3)
        title_width = Inches(12)
        title_height = Inches(1.0)
        
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_tf = title_box.text_frame
        title_tf.text = slide_data.get('title', 'No Title')
        title_tf.paragraphs[0].font.size = Pt(40)
        title_tf.paragraphs[0].font.bold = True

        # Content (Text) - Left side
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(6.0)
        height = Inches(5.0)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        # Handle content as either string or list
        content = slide_data.get('content', '')
        if isinstance(content, list):
            content_text = '\n'.join([f"• {item}" for item in content])
        else:
            content_text = content
            
        tf.text = content_text
        tf.word_wrap = True
        
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(20)

        # Image generation - Right side
        img_prompt = slide_data.get('image_prompt', '')
        img_filename = f"generated_images/slide_{i+1}.png"
        
        image_generated = False
        if generate_images and api_key and img_prompt:
            image_generated = generate_image_with_imagen(img_prompt, img_filename, api_key)
        
        if image_generated and os.path.exists(img_filename):
            # Add the generated image
            img_left = Inches(7.0)
            img_top = Inches(1.5)
            img_width = Inches(5.8)
            slide.shapes.add_picture(img_filename, img_left, img_top, width=img_width)
        else:
            # Create placeholder
            placeholder_left = Inches(7.0)
            placeholder_top = Inches(1.5)
            placeholder_width = Inches(5.8)
            placeholder_height = Inches(4.0)
            
            shape = slide.shapes.add_shape(
                1, # msoShapeRectangle
                placeholder_left, placeholder_top, placeholder_width, placeholder_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(238, 238, 238)
            shape.line.color.rgb = RGBColor(136, 136, 136)
            
            p_tf = shape.text_frame
            p_tf.text = f"Image Placeholder\n\nPrompt:\n{img_prompt}" if img_prompt else "No Image Prompt"
            p_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            p_tf.paragraphs[0].font.bold = True
            p_tf.paragraphs[0].font.size = Pt(14)
        
        # Add notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = f"Image Prompt: {img_prompt}"
        
        print(f"  ✅ Slide {i+1} completed\n")

    prs.save(output_file)
    print(f"\n✅ Presentation saved to {output_file}")
    print(f"📊 Created {len(slides_data)} slides")

if __name__ == "__main__":
    create_presentation(generate_images=True)
